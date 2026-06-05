"""Génère le support de soutenance (.pptx) · Urban Data Explorer.

RNCP40875 · Bloc 1 · Expert en Ingénierie des Données · Adam Beloucif & Émilien Morice.
Charte DSFR (Système de Design de l'État) : bleu France #000091,
rouge Marianne #E1000F, fond blanc, encadrés nets. La police Marianne
est rendue via sa substitution système officielle (Arial) pour un
affichage identique sur le poste du jury.

Dépendance : python-pptx (pip install python-pptx).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

# ── Palette DSFR ──────────────────────────────────────────────────────────
BLUE = RGBColor(0x00, 0x00, 0x91)        # bleu France
RED = RGBColor(0xE1, 0x00, 0x0F)         # rouge Marianne
INK = RGBColor(0x16, 0x16, 0x16)         # texte titre
GREY = RGBColor(0x3A, 0x3A, 0x3A)        # texte courant
MENTION = RGBColor(0x66, 0x66, 0x66)     # texte secondaire
LIGHT = RGBColor(0xF6, 0xF6, 0xF6)       # fond alt
CONTRAST = RGBColor(0xEE, 0xEE, 0xEE)    # fond contrasté
BORDER = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_TINT = RGBColor(0xE3, 0xE3, 0xFD)   # bleu France 925
BLUE_LIGHT = RGBColor(0x85, 0x85, 0xF6)  # bleu France clair (thème sombre)
EFREI_AZURE = RGBColor(0x14, 0x86, 0xC9)  # bleu eFrei (certificateur)
FONT = "Arial"
ASSETS = Path(__file__).resolve().parent / "assets"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.shadow.inherit = False
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    else:
        sp.line.fill.background()
    return sp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6):
    """runs: list of paragraphs ; chaque paragraphe = list de (texte, size, bold, color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb


def tricolore(slide, x, y, h):
    w = Inches(0.18)
    rect(slide, x, y, w, h, BLUE)
    rect(slide, x + w, y, w, h, WHITE, line=BORDER, line_w=Pt(0.5))
    rect(slide, x + 2 * w, y, w, h, RED)


def bullets(slide, x, y, w, h, items, size=15, gap=10, color=GREY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        # puce carrée bleu France
        rb = p.add_run(); rb.text = "▪  "; rb.font.name = FONT
        rb.font.size = Pt(size); rb.font.bold = True; rb.font.color.rgb = BLUE
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead; r1.font.name = FONT
            r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = INK
            r2 = p.add_run(); r2.text = rest; r2.font.name = FONT
            r2.font.size = Pt(size); r2.font.bold = False; r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = it; r.font.name = FONT
            r.font.size = Pt(size); r.font.bold = False; r.font.color.rgb = color
    return tb


# ── Slide builders ─────────────────────────────────────────────────────────

def content_header(slide, eyebrow, title):
    """Bandeau supérieur standard : liseré bleu, eyebrow, titre."""
    rect(slide, 0, 0, EMU_W, Inches(1.25), WHITE)
    rect(slide, Inches(0.6), Inches(0.34), Inches(0.09), Inches(0.62), BLUE)
    textbox(slide, Inches(0.85), Inches(0.3), Inches(11.8), Inches(0.7), [
        [(eyebrow.upper(), 12, True, MENTION)],
        [(title, 26, True, INK)],
    ], space_after=2)
    rect(slide, Inches(0.85), Inches(1.16), Inches(11.6), Pt(1.4), BORDER)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rect(s, 0, 0, EMU_W, EMU_H, WHITE)
    return s


def footer(slide, page):
    textbox(slide, Inches(0.85), Inches(7.05), Inches(8), Inches(0.3),
            [[("Urban Data Explorer · RNCP40875 Bloc 1 · Adam Beloucif & Émilien Morice", 9, False, MENTION)]])
    textbox(slide, Inches(11.5), Inches(7.05), Inches(1), Inches(0.3),
            [[(str(page), 9, True, BLUE)]], align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    # ── 1. Couverture ───────────────────────────────────────────────────
    s = new_slide(prs)
    tricolore(s, Inches(0.85), Inches(0.7), Inches(0.55))
    textbox(s, Inches(1.55), Inches(0.66), Inches(5), Inches(0.6), [
        [("RÉPUBLIQUE FRANÇAISE", 12, True, INK)],
        [("Données ouvertes · Ville de Paris", 10, False, MENTION)],
    ], space_after=1)
    textbox(s, Inches(0.85), Inches(2.5), Inches(11.6), Inches(1.8), [
        [("Urban Data Explorer", 52, True, BLUE)],
        [("Observatoire socio-urbain de Paris", 24, False, GREY)],
    ], space_after=8)
    rect(s, Inches(0.9), Inches(4.5), Inches(3.2), Pt(3), RED)
    textbox(s, Inches(0.85), Inches(4.8), Inches(11.6), Inches(1.6), [
        [("Architecture de stockage et de traitement de données", 16, True, INK)],
        [("RNCP40875 · Bloc 1 · Expert en Ingénierie des Données · Certificateur Efrei", 13, False, MENTION)],
        [("Adam Beloucif & Émilien Morice · Soutenance 2026", 13, False, MENTION)],
    ], space_after=5)
    # logo EFREI (certificateur) en haut a droite
    efrei = ASSETS / "efrei-logo.png"
    if efrei.exists():
        s.shapes.add_picture(str(efrei), Inches(10.25), Inches(0.6), width=Inches(2.5))
    # bandeau bas bicolore : bleu France (Etat) + azur EFREI (ecole)
    rect(s, 0, Inches(7.18), Inches(9.3), Inches(0.32), BLUE)
    rect(s, Inches(9.3), Inches(7.18), Inches(4.05), Inches(0.32), EFREI_AZURE)

    # ── 2. Sommaire ─────────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Plan", "Sommaire de la soutenance")
    left = [
        ("1. ", "Contexte & besoin de gouvernance"),
        ("2. ", "Architecture d'ensemble (médaillon)"),
        ("3. ", "Sources Open Data & géocodage"),
        ("4. ", "C1.1 · Base relationnelle & tests de charge"),
        ("5. ", "C1.2 · Base NoSQL Cassandra"),
        ("6. ", "C1.3 / C1.4 · Data Lake, scalabilité, résilience"),
    ]
    right = [
        ("7. ", "C2.1 · API REST FastAPI"),
        ("8. ", "C2.2 · Streaming Kafka temps réel"),
        ("9. ", "C2.3 / C2.4 · Transformation & optimisation"),
        ("10. ", "Indicateurs Gold & score de synthèse"),
        ("11. ", "Interface DSFR & accessibilité"),
        ("12. ", "Conformité grille & conclusion"),
    ]
    bullets(s, Inches(0.85), Inches(1.7), Inches(5.8), Inches(5), left, size=16, gap=14)
    bullets(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5), right, size=16, gap=14)
    footer(s, 2)

    # ── 3. Contexte ─────────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Contexte", "Besoin de gouvernance urbaine")
    bullets(s, Inches(0.85), Inches(1.7), Inches(7.0), Inches(5), [
        ("Problème métier : ", "les données urbaines de Paris (logement, mobilité, revenus, cadre de vie) sont dispersées sur des dizaines de jeux Open Data hétérogènes."),
        ("Besoin gouvernance : ", "disposer d'un observatoire unique, comparatif et cartographique pour éclairer les décisions d'aménagement par arrondissement."),
        ("Réponse : ", "une solution complète d'architecture de données · ingestion, stockage relationnel + NoSQL, Data Lake, API et restitution cartographique."),
        ("Cadre : ", "validation des compétences C1.1 à C2.4 du Bloc 1 (architecture de stockage et de traitement)."),
    ], size=15, gap=14)
    # carte chiffres clés
    rect(s, Inches(8.4), Inches(1.7), Inches(4.05), Inches(4.6), LIGHT, line=BORDER)
    rect(s, Inches(8.4), Inches(1.7), Inches(0.1), Inches(4.6), BLUE)
    textbox(s, Inches(8.75), Inches(2.0), Inches(3.5), Inches(4), [
        [("CHIFFRES CLÉS", 12, True, MENTION)],
        [("24", 34, True, BLUE)], [("sources de données intégrées", 12, False, GREY)],
        [("8", 34, True, BLUE)], [("familles thématiques", 12, False, GREY)],
        [("20", 34, True, BLUE)], [("arrondissements analysés", 12, False, GREY)],
        [("5", 34, True, BLUE)], [("niveaux de granularité", 12, False, GREY)],
    ], space_after=2)
    footer(s, 3)

    # ── 4. Architecture ─────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Architecture", "Pipeline médaillon de bout en bout")
    stages = [
        ("BRONZE", "Open Data brut\nParquet", BLUE_TINT, INK),
        ("SILVER", "Normalisé + géocodé\nParquet", RGBColor(0xB1,0xB1,0xF9), INK),
        ("GOLD", "Datamarts\nPostgreSQL", BLUE, WHITE),
    ]
    x = Inches(0.85)
    for name, desc, col, txt in stages:
        rect(s, x, Inches(1.9), Inches(2.7), Inches(1.6), col, line=BORDER)
        textbox(s, x + Inches(0.15), Inches(2.05), Inches(2.4), Inches(1.3), [
            [(name, 18, True, txt)],
            [(desc, 12, False, txt)],
        ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=4)
        x += Inches(3.05)
    textbox(s, Inches(0.85), Inches(3.65), Inches(8.4), Inches(0.4),
            [[("Polars (moteur Rust) →  jointures spatiales IRIS  →  schéma en étoile", 13, True, GREY)]],
            align=PP_ALIGN.CENTER)
    bullets(s, Inches(0.85), Inches(4.35), Inches(11.6), Inches(2.4), [
        ("Ingestion : ", "ETL batch Polars (Bronze→Silver→Gold) + collecte temps réel Kafka."),
        ("Stockage : ", "PostgreSQL (relationnel étoile), Cassandra (NoSQL événements), Data Lake HDFS/Parquet."),
        ("Exposition : ", "API REST FastAPI documentée (Swagger /docs) → frontend React/TypeScript MapLibre (2D IGN) + Mapbox (3D)."),
        ("Orchestration : ", "Docker Compose multi-services (11 conteneurs, profils streaming & lake)."),
    ], size=14, gap=10)
    footer(s, 4)

    # ── 5. Sources & géocodage ──────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Données", "Sources Open Data & géocodage")
    bullets(s, Inches(0.85), Inches(1.7), Inches(7.0), Inches(5), [
        ("24 sources, 8 familles : ", "mobilité, éducation, logement, culture, santé, espaces verts, services publics, pression urbaine."),
        ("Sources réelles : ", "DVF 2023 (prix immobilier) et INSEE Filosofi 2020 (revenus) alimentent le cœur logement."),
        ("Géocodage offline : ", "point-in-polygon (ray-casting) sur les contours IRIS, sans appel réseau ; fallback API adresse.data.gouv.fr."),
        ("Qualité tracée : ", "drapeaux quality_flag (rattachement) et data_source (réel / référence)."),
    ], size=14, gap=12)
    rect(s, Inches(8.4), Inches(1.7), Inches(4.05), Inches(4.6), LIGHT, line=BORDER)
    rect(s, Inches(8.4), Inches(1.7), Inches(0.1), Inches(4.6), RED)
    textbox(s, Inches(8.75), Inches(2.0), Inches(3.5), Inches(4.1), [
        [("RÉPARTITION PAR FAMILLE", 11, True, MENTION)],
        [("Mobilité", 13, True, INK), ("   7", 13, False, GREY)],
        [("Logement", 13, True, INK), ("   3", 13, False, GREY)],
        [("Éducation", 13, True, INK), ("   3", 13, False, GREY)],
        [("Espaces verts", 13, True, INK), ("   3", 13, False, GREY)],
        [("Services publics", 13, True, INK), ("   3", 13, False, GREY)],
        [("Culture", 13, True, INK), ("   2", 13, False, GREY)],
        [("Santé", 13, True, INK), ("   2", 13, False, GREY)],
        [("Pression urbaine", 13, True, INK), ("   1", 13, False, GREY)],
    ], space_after=6)
    footer(s, 5)

    # ── 6. C1.1 Relationnel ─────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C1.1 · Base relationnelle", "Schéma en étoile & tests de charge")
    bullets(s, Inches(0.85), Inches(1.7), Inches(6.0), Inches(3), [
        ("Modèle en étoile : ", "fact_arrondissement_dashboard relié à dim_arrondissement / dim_iris."),
        ("Intégrité : ", "clés primaires, clés étrangères, contraintes NOT NULL, normalisation."),
        ("Index de performance : ", "sur clés étrangères et indices d'attractivité (Index Scan)."),
        ("PostgreSQL 15 ", "(conteneur Docker)."),
    ], size=14, gap=9)
    # tableau bench
    rows = [
        ("Métrique (bench concurrent)", "Valeur"),
        ("Requêtes réussies", "500 / 500"),
        ("Débit moyen", "1 190 req/s"),
        ("Latence moyenne", "1,80 ms"),
        ("Latence p95", "3,50 ms"),
        ("Intégrité référentielle", "3/3 PASSED"),
    ]
    add_table(s, Inches(7.1), Inches(1.8), Inches(5.35), rows, col0_w=3.4)
    footer(s, 6)

    # ── 7. C1.2 NoSQL ───────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C1.2 · Base NoSQL", "Cassandra pour les événements de streaming")
    bullets(s, Inches(0.85), Inches(1.7), Inches(7.0), Inches(5), [
        ("Choix justifié : ", "Cassandra pour des écritures massives, append-only, orientées requêtes · adapté aux flux d'événements semi-structurés."),
        ("Modélisation orientée requêtes : ", "table events_by_type partitionnée par event_type."),
        ("Tri chronologique : ", "CLUSTERING ORDER BY event_time DESC pour servir les événements récents."),
        ("Cycle de vie : ", "TTL de 604 800 s (7 jours) · purge automatique des snapshots."),
        ("Complémentarité SQL/NoSQL : ", "PostgreSQL pour l'analytique structurée, Cassandra pour le flux temps réel."),
    ], size=14, gap=11)
    rect(s, Inches(8.4), Inches(1.7), Inches(4.05), Inches(4.4), INK)
    textbox(s, Inches(8.7), Inches(1.95), Inches(3.5), Inches(4), [
        [("events_by_type", 13, True, BLUE_LIGHT)],
        [("PRIMARY KEY", 11, True, WHITE)],
        [("(event_type,", 12, False, CONTRAST)],
        [(" event_time DESC,", 12, False, CONTRAST)],
        [(" event_id)", 12, False, CONTRAST)],
        [("TTL 7 jours", 12, True, RED)],
        [("RF = 1 · SimpleStrategy", 11, False, MENTION)],
    ], space_after=6)
    footer(s, 7)

    # ── 8. C1.3 / C1.4 ──────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C1.3 / C1.4", "Data Lake, scalabilité & résilience")
    bullets(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(2.4), [
        ("C1.3 · Data Lake : ", "stockage multivarié Bronze/Silver/Gold en Parquet (colonnaire, compressé) sur HDFS ; intégration de sources variées."),
        ("Indicateurs suivis : ", "latence, débit, volumétrie · mesurés sur le benchmark PostgreSQL et le pipeline."),
        ("Sécurité & protection : ", "secrets via variables d'environnement (.env gitignored), CORS maîtrisé, données publiques anonymes."),
    ], size=14, gap=10)
    bullets(s, Inches(0.85), Inches(4.3), Inches(11.6), Inches(2.4), [
        ("C1.4 · Scalabilité : ", "architecture conteneurisée Docker Compose (11 services), profils activables (streaming, lake)."),
        ("Haute disponibilité : ", "Zookeeper + Kafka, Hadoop namenode/datanode, Hive metastore."),
        ("Résilience : ", "mode local-first hors ligne (Parquet + fallback) garantissant la continuité même sans cluster."),
    ], size=14, gap=10)
    footer(s, 8)

    # ── 9. C2.1 API ─────────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C2.1 · API REST", "FastAPI : exposition des données")
    bullets(s, Inches(0.85), Inches(1.7), Inches(6.5), Inches(5), [
        ("Technologie adaptée : ", "FastAPI (ASGI, Pydantic v2), documentation auto (Swagger /docs)."),
        ("Routeurs métier : ", "health, auth, catalog, datamarts, pipeline, events, repo."),
        ("Authentification JWT : ", "OAuth2, rôles viewer / admin, jeton signé HS256 (60 min)."),
        ("Autorisations & quotas : ", "route admin protégée ; limite par IP (anonyme 120, connecté 600 req/min → 429)."),
        ("CORS restreint : ", "à l'origine du frontend (variable UDE_CORS_ORIGINS)."),
    ], size=13, gap=10)
    # capture reelle du Swagger /docs
    api_img = ROOT / "ude-api.png"
    if api_img.exists():
        rect(s, Inches(7.55), Inches(1.62), Inches(5.0), Inches(4.86), WHITE, line=BORDER)
        s.shapes.add_picture(str(api_img), Inches(7.62), Inches(1.69), height=Inches(4.72))
    footer(s, 9)

    # ── 10. C2.2 Streaming ──────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C2.2 · Streaming", "Kafka : ingestion temps réel")
    bullets(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(5), [
        ("Veille & choix : ", "Apache Kafka retenu pour le découplage producteur/consommateur et la tolérance aux pannes."),
        ("Producteur : ", "publie les flux urbains (Vélib disponibilité, chantiers perturbants) au fil de l'eau."),
        ("Consommateur : ", "lit le topic et écrit les snapshots dans Cassandra (events_by_type)."),
        ("Micro-batch & temps réel : ", "traitement au fur et à mesure de la disponibilité, restitué dans le flux d'événements de l'interface."),
        ("Infrastructure : ", "Zookeeper + Kafka orchestrés via le profil 'streaming' de Docker Compose."),
    ], size=15, gap=13)
    footer(s, 10)

    # ── 11. C2.3 / C2.4 ─────────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "C2.3 / C2.4", "Transformation & optimisation des pipelines")
    bullets(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(2.4), [
        ("C2.3 · Intégration : ", "Polars normalise les codes (75101→75001), parse les geo_point, fusionne des sources hétérogènes."),
        ("Jointures spatiales : ", "rattachement IRIS par point-in-polygon, puis agrégation multidimensionnelle par arrondissement."),
        ("Structuration analytique : ", "datamarts Gold (dashboard + timeline 12 mois) prêts pour l'analyse."),
    ], size=14, gap=10)
    bullets(s, Inches(0.85), Inches(4.3), Inches(11.6), Inches(2.4), [
        ("C2.4 · Performance : ", "Parquet colonnaire compressé (Silver/Gold), moteur Polars vectorisé (Rust)."),
        ("Accès rapide : ", "index PostgreSQL → latence p95 < 4 ms sous charge concurrente (10 threads)."),
        ("Indicateurs de transfert : ", "débit ~1 190 req/s mesuré, temps de run du pipeline suivi."),
        ("Qualité : ", "quality_flag garantissant des analyses exploitables."),
    ], size=14, gap=10)
    footer(s, 11)

    # ── 12. Indicateurs Gold ────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Restitution", "Indicateurs Gold & score de synthèse")
    cards = [
        ("Marché immobilier", "Prix au m² réel & ventes (DVF 2023)"),
        ("Logement social", "Part (%) & nombre de logements"),
        ("Revenus", "Revenu médian réel (INSEE Filosofi)"),
        ("Accessibilité", "m² achetables / an de revenu (prix ÷ revenu)"),
        ("Cadre de vie & env.", "Transports, équipements, espaces verts"),
        ("Score global", "Moyenne des indices /100"),
    ]
    x0, y0 = Inches(0.85), Inches(1.8)
    cw, ch = Inches(3.75), Inches(1.7)
    for i, (t, d) in enumerate(cards):
        cx = x0 + (i % 3) * (cw + Inches(0.18))
        cy = y0 + (i // 3) * (ch + Inches(0.22))
        last = (i == 5)
        rect(s, cx, cy, cw, ch, BLUE if last else WHITE, line=BORDER)
        rect(s, cx, cy, Inches(0.08), ch, RED if last else BLUE)
        textbox(s, cx + Inches(0.22), cy + Inches(0.2), cw - Inches(0.35), ch - Inches(0.3), [
            [(t, 16, True, WHITE if last else INK)],
            [(d, 12, False, CONTRAST if last else GREY)],
        ], space_after=5)
    footer(s, 12)

    # ── 13. Interface clair ─────────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Interface", "Dashboard cartographique aux codes DSFR")
    img = ROOT / "ude-light.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0.85), Inches(1.7), height=Inches(4.7))
    bullets(s, Inches(8.7), Inches(1.8), Inches(3.9), Inches(4.6), [
        ("Charte DSFR : ", "Marianne, bleu France, rouge Marianne."),
        ("Fond IGN : ", "tuiles Géoplateforme officielles."),
        ("5 granularités : ", "ville → bâtiment."),
        ("Choroplèthe : ", "échelle séquentielle bleu France."),
        ("Comparatif : ", "2 arrondissements côte à côte."),
    ], size=13, gap=10)
    footer(s, 13)

    # ── 14. Interface sombre + a11y ─────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Interface", "Thème sombre & accessibilité")
    img = ROOT / "ude-dark.png"
    if img.exists():
        s.shapes.add_picture(str(img), Inches(0.85), Inches(1.7), height=Inches(4.7))
    bullets(s, Inches(8.7), Inches(1.8), Inches(3.9), Inches(4.6), [
        ("Thème sombre DSFR : ", "bascule clair/sombre persistante."),
        ("Contraste : ", "conforme WCAG AA (rouge texte #CE0500)."),
        ("Focus visibles ", "& labels ARIA sur les contrôles."),
        ("Mouvement réduit : ", "prefers-reduced-motion respecté."),
        ("Chiffres tabulaires ", "pour les données alignées."),
    ], size=13, gap=10)
    footer(s, 14)

    # ── 15. Grille de conformité ────────────────────────────────────────
    s = new_slide(prs)
    content_header(s, "Conformité", "Couverture de la grille RNCP40875 · Bloc 1")
    rows = [
        ("Compétence", "Preuve dans le projet"),
        ("C1.1 Relationnel", "Schéma étoile PostgreSQL + tests de charge"),
        ("C1.2 NoSQL", "Cassandra events_by_type, TTL 7 j"),
        ("C1.3 Data Lake", "Bronze/Silver/Gold Parquet sur HDFS"),
        ("C1.4 Scalabilité", "Docker Compose 11 services, local-first"),
        ("C2.1 API", "FastAPI /docs + auth JWT (rôles) + quotas par IP"),
        ("C2.2 Streaming", "Producteur/consommateur Kafka"),
        ("C2.3 Transformation", "Polars + sources réelles DVF / INSEE Filosofi"),
        ("C2.4 Optimisation", "Parquet, index, p95 < 4 ms, qualité tracée"),
    ]
    add_table(s, Inches(0.85), Inches(1.7), Inches(11.6), rows, col0_w=3.3, fs=13)
    footer(s, 15)

    # ── 16. Conclusion ──────────────────────────────────────────────────
    s = new_slide(prs)
    rect(s, 0, 0, EMU_W, EMU_H, BLUE)
    tricolore(s, Inches(0.85), Inches(0.7), Inches(0.5))
    textbox(s, Inches(0.85), Inches(2.2), Inches(11.6), Inches(2), [
        [("Conclusion & perspectives", 38, True, WHITE)],
    ])
    bullets(s, Inches(0.9), Inches(3.4), Inches(11.4), Inches(3), [
        ("Solution complète : ", "ingestion → stockage SQL/NoSQL/Lake → API → restitution cartographique, couvrant C1.1 à C2.4."),
        ("Robuste : ", "fonctionne en local-first hors ligne comme en mode distribué conteneurisé."),
        ("Soignée : ", "interface aux standards de l'État (DSFR), accessible et bilingue données live / offline."),
        ("Perspectives : ", "authentification JWT active, branchement DVF/INSEE réels, déploiement cloud, alerting."),
    ], size=15, gap=12, color=CONTRAST)
    textbox(s, Inches(0.9), Inches(6.7), Inches(8), Inches(0.5),
            [[("Merci de votre attention · questions bienvenues.", 14, True, WHITE)]])
    # logo EFREI blanc en bas a droite
    efrei_w = ASSETS / "efrei-logo-white.png"
    if efrei_w.exists():
        s.shapes.add_picture(str(efrei_w), Inches(10.6), Inches(6.45), width=Inches(2.2))

    # ── Notes de l'orateur · minutage pour 10 minutes ───────────────────
    # (annexe = diapos de reserve pour les questions, non comptees dans les 10 min)
    NOTES = [
        "[~0:30] Accroche : Paris, le logement, des donnees dispersees. On a construit l'observatoire qui les reunit.",
        "[~0:20] Plan rapide. Annoncer la demo live a la fin.",
        "[~1:00] Probleme metier + chiffres cles (24 sources, 20 arrondissements). Pourquoi une gouvernance en a besoin.",
        "[~1:00] Le voyage de la donnee : Bronze/Silver/Gold, Polars, puis API et carte. Insister sur le bout-en-bout.",
        "[~1:00] Sources : 24, equilibrees. Surtout : prix DVF et revenus INSEE REELS. Geocodage hors ligne.",
        "[~0:40] C1.1 : etoile PostgreSQL + tests de charge (1190 req/s, p95 3,5 ms).",
        "[~0:30] C1.2 : Cassandra pour le flux, tri par date, purge 7 jours.",
        "[~0:40] C1.3/C1.4 : Data Lake Parquet, Docker, healthchecks, volumes persistants, local-first.",
        "[~0:50] C2.1 : API FastAPI, auth JWT + quotas par IP. Montrer /docs en demo.",
        "[~0:30] C2.2 : Kafka temps reel vers Cassandra.",
        "[~0:40] C2.3/C2.4 : fusion sources reelles, Parquet, index, qualite tracee.",
        "[~0:40] Indicateurs Gold dont l'accessibilite prix/revenus (KPI demande).",
        "[~1:00] DEMO LIVE : carte, choix d'un arrondissement, comparaison, theme sombre.",
        "[Annexe] Accessibilite / theme sombre si question.",
        "[~0:30] Grille : on coche C1.1 a C2.4.",
        "[~0:30] Conclusion : solution complete, reelle, soignee. Merci. Questions.",
    ]
    for slide, note in zip(prs.slides, NOTES):
        slide.notes_slide.notes_text_frame.text = note

    out = ROOT / "soutenance.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = ROOT / "soutenance-maj.pptx"
        prs.save(str(out))
        print("ATTENTION : soutenance.pptx etait ouvert/verrouille -> ecriture dans soutenance-maj.pptx")
    print(f"Support genere: {out}  ({len(prs.slides._sldIdLst)} slides)")


def add_table(slide, x, y, w, rows, col0_w=3.0, fs=13):
    """Tableau DSFR à 2 colonnes ; rows[0] = en-tête."""
    from pptx.util import Inches as _In
    n = len(rows)
    row_h = _In(0.46)
    tbl_shape = slide.shapes.add_table(n, 2, x, y, w, row_h * n)
    table = tbl_shape.table
    table.columns[0].width = _In(col0_w)
    table.columns[1].width = w - _In(col0_w)
    for r, (a, b) in enumerate(rows):
        for c, val in enumerate((a, b)):
            cell = table.cell(r, c)
            cell.margin_left = Pt(8); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.name = FONT
            run.font.size = Pt(fs)
            run.font.bold = (r == 0) or (c == 0)
            run.font.color.rgb = WHITE if r == 0 else (BLUE if c == 0 else GREY)


if __name__ == "__main__":
    build()
