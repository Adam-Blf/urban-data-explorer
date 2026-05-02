"""Génère le support de soutenance Urban Data Explorer.

Sortie · docs/Soutenance_Urban_Data_Explorer.pptx
Thème sombre premium · EFREI logo en footer.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
OUT = DOCS / "Soutenance_Urban_Data_Explorer.pptx"

NAVY  = RGBColor(0x0B, 0x10, 0x20)
NAVY2 = RGBColor(0x11, 0x17, 0x2D)
NAVY3 = RGBColor(0x1A, 0x21, 0x42)
WHITE = RGBColor(0xF4, 0xF6, 0xFB)
GREY  = RGBColor(0xAA, 0xB2, 0xCC)
GOLD  = RGBColor(0xFF, 0xD1, 0x66)
TEAL  = RGBColor(0x06, 0xD6, 0xA0)
PINK  = RGBColor(0xEF, 0x47, 0x6F)
BLUE  = RGBColor(0x11, 0x8A, 0xB2)

WIDTH_IN, HEIGHT_IN = 13.333, 7.5  # 16:9 widescreen


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def fill_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(WIDTH_IN), Inches(HEIGHT_IN))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(
    slide, text, left, top, width, height,
    *, size=18, bold=False, color=WHITE, font="Helvetica",
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def add_round_rect(slide, left, top, width, height, color, radius=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_logo(slide, left, top, height_in, white=True):
    """Logo blanc par défaut (mieux contrasté sur fond sombre du deck)."""
    logo = ASSETS / ("efrei-logo-white.png" if white else "efrei-logo.png")
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(left), Inches(top), height=Inches(height_in))


def add_footer(slide, page_num, total):
    add_rect(slide, 0, HEIGHT_IN - 0.35, WIDTH_IN, 0.35, NAVY2)
    add_text(slide, "Urban Data Explorer · Adam Beloucif · Emilien Morice",
             0.5, HEIGHT_IN - 0.32, 8, 0.3,
             size=9, color=GREY)
    add_text(slide, f"{page_num} / {total}",
             WIDTH_IN - 1.5, HEIGHT_IN - 0.32, 1, 0.3,
             size=9, color=GREY, align=PP_ALIGN.RIGHT)
    add_logo(slide, WIDTH_IN - 1.0, HEIGHT_IN - 0.30, 0.25)


def add_kicker(slide, text, color=GOLD):
    add_text(slide, text.upper(), 0.7, 0.55, 6, 0.4,
             size=11, bold=True, color=color)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)

    # accent bar
    add_round_rect(slide, 0.7, 1.6, 0.18, 1.6, GOLD, radius=0.5)

    add_text(slide, "Urban Data Explorer", 1.1, 1.5, 11, 1.4,
             size=58, bold=True, color=WHITE, line_spacing=1.0)
    add_text(slide,
             "Plateforme open data Paris · pipeline médaillon, API REST sécurisée\net dashboard cartographique du logement parisien.",
             1.1, 3.05, 11, 1.4,
             size=20, color=GREY, line_spacing=1.3)

    add_text(slide, "ADAM BELOUCIF",  1.1, 5.0, 5, 0.4, size=14, bold=True, color=GOLD)
    add_text(slide, "EMILIEN MORICE", 1.1, 5.4, 5, 0.4, size=14, bold=True, color=TEAL)

    add_text(slide,
             "M1 Data Engineering & IA · EFREI Paris · Architecture de Données · 2026",
             1.1, 5.9, 11, 0.4, size=13, color=GREY)

    # big logo top-right
    add_logo(slide, WIDTH_IN - 2.6, 0.55, 1.0)
    return slide


def slide_problem(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "01 · Le problème", PINK)

    add_text(slide, "Paris, marché du logement\nfragmenté, illisible.",
             0.7, 1.0, 11, 2.0, size=42, bold=True, line_spacing=1.05)

    # 3 numéros clés
    cols = [
        ("> 10 000 €", "prix m² médian intra-muros", GOLD),
        ("17 sources", "données ouvertes hétérogènes", TEAL),
        ("20 arr.", "réalités économiques différentes", PINK),
    ]
    for i, (big, label, color) in enumerate(cols):
        x = 0.7 + i * 4.1
        add_round_rect(slide, x, 3.6, 3.8, 1.6, NAVY2)
        add_text(slide, big,   x + 0.3, 3.75, 3.5, 0.8, size=34, bold=True, color=color)
        add_text(slide, label, x + 0.3, 4.55, 3.5, 0.6, size=12, color=GREY)

    add_text(slide,
             "Comment lire le marché du logement à l'échelle de l'arrondissement\nen un seul écran, croisé avec revenus, équipements, environnement ?",
             0.7, 5.6, 12, 1.2, size=16, color=WHITE, line_spacing=1.3)

    add_footer(slide, n, total)


def slide_answer(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "02 · La réponse", GOLD)

    add_text(slide, "Médaillon · API · Dashboard.",
             0.7, 1.0, 12, 1.6, size=54, bold=True, line_spacing=1.0)
    add_text(slide,
             "Une plateforme complète, du pipeline d'ingestion à la carto interactive.",
             0.7, 2.7, 12, 0.6, size=18, color=GREY)

    cards = [
        ("17",  "sources de données",         "data.gouv · OpenData Paris · INSEE · Airparif", GOLD),
        ("4",   "indicateurs composites",     "accessibilité · tension · effort social · attractivité", TEAL),
        ("4",   "niveaux cartographiques",    "choroplèthe · POI · timeline · comparaison",    PINK),
        ("19",  "commits granulaires",        "alternance Adam / Emilien · doc + ADRs",        BLUE),
    ]
    for i, (big, mid, sub, color) in enumerate(cards):
        x = 0.7 + (i % 4) * 3.1
        add_round_rect(slide, x, 4.0, 2.9, 2.3, NAVY2)
        add_text(slide, big, x + 0.3, 4.15, 2.6, 1.0, size=44, bold=True, color=color)
        add_text(slide, mid, x + 0.3, 5.05, 2.6, 0.4, size=13, bold=True, color=WHITE)
        add_text(slide, sub, x + 0.3, 5.45, 2.6, 0.8, size=10, color=GREY, line_spacing=1.3)

    add_footer(slide, n, total)


def slide_architecture(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "03 · Architecture", GOLD)
    add_text(slide, "Pipeline médaillon · Bronze → Silver → Gold",
             0.7, 1.0, 12, 0.7, size=30, bold=True)

    layers = [
        ("BRONZE /raw",  "Snapshot brut versionné Parquet partitionné year/month/day",       "feeder.py · CLI paramétrable",                 GOLD),
        ("SILVER /silver","Cleaning + 6 règles validation + jointure spatiale + windows",    "processor.py · Polars + Shapely + cache",       TEAL),
        ("GOLD /gold",   "DuckDB relationnel · dim, facts, KPI, timeline",                   "datamart.py · 4 indicateurs composites",        PINK),
        ("API",          "FastAPI 0.115 · JWT · pagination · filtres typés",                 "uvicorn :8000 · OpenAPI auto",                  BLUE),
        ("FRONTEND",     "MapLibre + Chart.js · 4 niveaux d'info · timeline · comparaison",  "ESM via CDN · zéro build · déployable Vercel",  WHITE),
    ]
    y = 2.0
    for tag, desc, sub, color in layers:
        add_round_rect(slide, 0.7, y, 1.9, 0.85, color)
        add_text(slide, tag, 0.85, y + 0.2, 1.7, 0.4, size=13, bold=True, color=NAVY)
        add_text(slide, desc, 2.8, y + 0.05, 9.7, 0.5, size=14, bold=True, color=WHITE)
        add_text(slide, sub,  2.8, y + 0.45, 9.7, 0.5, size=11, color=GREY)
        y += 1.0

    add_footer(slide, n, total)


def slide_sources(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "04 · Sources de données", GOLD)
    add_text(slide, "17 datasets ouverts, > 350 000 enregistrements.",
             0.7, 1.0, 12, 0.7, size=26, bold=True)

    rows = [
        ("Cœur",          "DVF géolocalisé",            "data.gouv (DGFiP)",     "~250 k"),
        ("Cœur",          "Logements sociaux",          "OpenData Paris",        "~5 k"),
        ("Cœur",          "Filosofi 2021",              "INSEE",                 "1 / commune"),
        ("Cœur",          "Arrondissements GeoJSON",    "OpenData Paris",        "20 polygones"),
        ("Environnement", "Qualité air Airparif",       "OpenData Paris",        "~30 k"),
        ("Transport",     "Vélib stations",             "OpenData Paris",        "~1 400"),
        ("Transport",     "Belib bornes IRVE",          "OpenData Paris",        "~600"),
        ("Transport",     "Aménagements cyclables",     "OpenData Paris",        "~12 000"),
        ("Service public","Écoles élémentaires",        "OpenData Paris",        "~330"),
        ("Service public","Collèges",                   "OpenData Paris",        "~120"),
        ("Service public","Sanisettes",                 "OpenData Paris",        "~750"),
        ("Commerce",      "Marchés découverts",         "OpenData Paris",        "~80"),
        ("Culture",       "Que faire à Paris",          "OpenData Paris",        "~10 000 / an"),
        ("Culture",       "Musées de France",           "data.gouv",             "~150 IDF"),
        ("Culture",       "Monuments historiques",      "data.gouv",             "~3 800 IDF"),
        ("Santé",         "Hôpitaux IDF",               "data.gouv",             "~370"),
        ("Environnement", "Espaces verts publics",      "OpenData Paris",        "~3 000"),
    ]
    headers = ("Catégorie", "Dataset", "Provider", "Volume")
    col_x = [0.7, 2.8, 6.4, 10.4]
    col_w = [2.0, 3.5, 3.9, 2.5]

    add_round_rect(slide, 0.5, 1.85, 12.4, 0.5, NAVY3)
    for h, x, w in zip(headers, col_x, col_w):
        add_text(slide, h.upper(), x, 1.92, w, 0.4, size=10, bold=True, color=GOLD)

    y = 2.4
    row_h = 0.27
    for i, row in enumerate(rows):
        if i % 2 == 0:
            add_rect(slide, 0.5, y - 0.02, 12.4, row_h + 0.02, NAVY2)
        for cell, x, w in zip(row, col_x, col_w):
            color = TEAL if cell.startswith("~") or "/" in cell else WHITE
            add_text(slide, cell, x, y, w, row_h, size=9, color=color)
        y += row_h

    add_footer(slide, n, total)


def slide_indicators(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "05 · 4 indicateurs composites", GOLD)
    add_text(slide, "Force de proposition · z-scores et formules métier.",
             0.7, 1.0, 12, 0.7, size=22, bold=True)

    cards = [
        ("idx_accessibilite", "revenu_median / (prix_m2 × 50)",
         "> 0.6 = accessible\n< 0.3 = très tendu", GOLD),
        ("idx_tension", "transactions_an × 1000 / parc_logements",
         "rotation pour 1 000\nlogements", TEAL),
        ("idx_effort_social", "log_sociaux × 10 000 / population",
         "unités sociales pour\n10 000 habitants", PINK),
        ("idx_attractivite", "0.30·z(access) + 0.40·z(POI) + 0.15·z(effort) − 0.15·z(tension)",
         "z-score composite\n0 = moyenne Paris", BLUE),
    ]
    for i, (name, formula, lecture, color) in enumerate(cards):
        x = 0.7 + (i % 2) * 6.2
        y = 2.0 + (i // 2) * 2.5
        add_round_rect(slide, x, y, 5.9, 2.3, NAVY2)
        add_round_rect(slide, x, y, 0.15, 2.3, color)
        add_text(slide, name, x + 0.3, y + 0.2, 5.5, 0.4, size=15, bold=True, color=color, font="Courier New")
        add_text(slide, "FORMULE", x + 0.3, y + 0.65, 5.5, 0.3, size=9, bold=True, color=GREY)
        add_text(slide, formula, x + 0.3, y + 0.95, 5.5, 0.5, size=12, color=WHITE, font="Courier New")
        add_text(slide, "LECTURE", x + 0.3, y + 1.45, 5.5, 0.3, size=9, bold=True, color=GREY)
        add_text(slide, lecture, x + 0.3, y + 1.75, 5.5, 0.5, size=12, color=WHITE, line_spacing=1.2)

    add_footer(slide, n, total)


def slide_pipeline(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "06 · Pipeline médaillon", GOLD)
    add_text(slide, "Idempotent · paramétrable · loggué.",
             0.7, 1.0, 12, 0.7, size=26, bold=True)

    bullets = [
        ("Partitionnement",     "year=YYYY/month=MM/day=DD sur Bronze + Silver"),
        ("CLI sans hardcode",   "Settings pydantic depuis .env · 0 chemin codé en dur"),
        ("Validation",          "6 règles métier sur DVF · log des rejets par règle"),
        ("Jointure spatiale",   "Shapely point-in-polygon sur GeoJSON arrondissements"),
        ("Window functions",    "MEDIAN OVER · RANK OVER · LAG OVER pour delta MoM"),
        ("Cache",               "Polars cache() sur Silver DVF · équivalent persist Spark"),
    ]
    for i, (label, desc) in enumerate(bullets):
        y = 2.0 + i * 0.75
        add_round_rect(slide, 0.7, y, 0.4, 0.4, GOLD)
        add_text(slide, str(i + 1), 0.7, y + 0.05, 0.4, 0.3, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, label, 1.3, y + 0.02, 4, 0.4, size=15, bold=True, color=WHITE)
        add_text(slide, desc,  1.3, y + 0.4,  10, 0.35, size=12, color=GREY)

    # code snippet card
    add_round_rect(slide, 8.3, 2.0, 4.4, 4.2, NAVY2)
    add_text(slide, "EXEMPLE", 8.5, 2.15, 4, 0.3, size=10, bold=True, color=GOLD)
    code = ("$ python -m pipeline.feeder \\\n"
            "      --source dvf --year 2024\n"
            "$ python -m pipeline.processor \\\n"
            "      --source dvf\n"
            "$ python -m pipeline.datamart \\\n"
            "      --build all\n\n"
            "→ /raw/dvf/year=2024/...\n"
            "→ /silver/dvf/...\n"
            "→ /gold/urban.duckdb")
    add_text(slide, code, 8.5, 2.5, 4, 3.6, size=11, color=TEAL, font="Courier New", line_spacing=1.3)

    add_footer(slide, n, total)


def slide_api(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "07 · API REST FastAPI", GOLD)
    add_text(slide, "JWT · pagination · filtres · OpenAPI auto.",
             0.7, 1.0, 12, 0.7, size=26, bold=True)

    eps = [
        ("POST", "/auth/login",                        "username/password → JWT", PINK),
        ("GET",  "/health",                            "liveness + version",       GREY),
        ("GET",  "/datamarts/arrondissements",         "KPI · pagination · filtres", TEAL),
        ("GET",  "/datamarts/timeline",                "série mensuelle prix m²",  TEAL),
        ("GET",  "/datamarts/indicators",              "4 indicateurs pivotés",     TEAL),
        ("GET",  "/geo/arrondissements.geojson",       "fond enrichi des KPI",      GOLD),
        ("GET",  "/geo/poi.geojson",                   "POI filtrables par catégorie", GOLD),
    ]
    y = 2.0
    for verb, path, desc, color in eps:
        add_round_rect(slide, 0.7, y, 0.7, 0.45, color)
        add_text(slide, verb, 0.75, y + 0.07, 0.65, 0.35, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, path, 1.5, y + 0.05, 6.5, 0.4, size=14, bold=True, color=WHITE, font="Courier New")
        add_text(slide, desc, 8.0, y + 0.05, 5, 0.4, size=12, color=GREY)
        y += 0.65

    add_footer(slide, n, total)


def slide_dashboard(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "08 · Dashboard", GOLD)
    add_text(slide, "MapLibre + Chart.js · 4 niveaux d'information.",
             0.7, 1.0, 12, 0.7, size=26, bold=True)

    layers = [
        ("01", "Choroplèthe par arrondissement", "Indicateur sélectionnable · interpolation continue", GOLD),
        ("02", "Couches POI activables",         "Transport · services publics · commerce · culture · santé · environnement", TEAL),
        ("03", "Timeline mensuelle",             "Slider qui repaint la carte en direct (2019-2024)", PINK),
        ("04", "Mode comparaison",               "Radar 5 dimensions entre 2 arrondissements", BLUE),
    ]
    for i, (num, title, desc, color) in enumerate(layers):
        y = 2.1 + i * 1.2
        add_round_rect(slide, 0.7, y, 12.1, 1.0, NAVY2)
        add_text(slide, num, 0.9, y + 0.15, 1, 0.7, size=42, bold=True, color=color, font="Helvetica")
        add_text(slide, title, 2.2, y + 0.1, 9, 0.5, size=18, bold=True, color=WHITE)
        add_text(slide, desc,  2.2, y + 0.55, 10, 0.4, size=12, color=GREY)

    add_footer(slide, n, total)


def slide_choices(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "09 · Choix techniques", GOLD)
    add_text(slide, "5 ADRs documentés.", 0.7, 1.0, 12, 0.7, size=26, bold=True)

    items = [
        ("ADR-001", "Polars + DuckDB > Spark",
         "Volume cible (~250 k DVF + 50 k POI) tient en RAM · 5× plus léger, démarrage instant.", GOLD),
        ("ADR-002", "Partitionnement year/month/day",
         "Idempotence stricte · replays sans doublons.", TEAL),
        ("ADR-003", "Indicateurs en z-score",
         "Comparable, lisible, robuste à l'évolution mensuelle.", PINK),
        ("ADR-004", "DuckDB read-only par requête",
         "Pas de pool · lecture concurrente sécurisée.", BLUE),
        ("ADR-005", "Frontend statique CDN-only",
         "Démo en 5 min sur n'importe quel host · pas de bundler.", WHITE),
    ]
    y = 2.0
    for tag, title, desc, color in items:
        add_round_rect(slide, 0.7, y, 1.3, 0.85, color)
        add_text(slide, tag, 0.9, y + 0.25, 1, 0.4, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, title, 2.2, y + 0.05, 10, 0.5, size=16, bold=True, color=WHITE)
        add_text(slide, desc,  2.2, y + 0.45, 10, 0.4, size=12, color=GREY)
        y += 1.0

    add_footer(slide, n, total)


def slide_demo(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)

    # vertical accent
    add_round_rect(slide, 0.7, 2.0, 0.18, 3.5, TEAL, radius=0.5)

    add_text(slide, "Démonstration", 1.1, 2.2, 11, 1.4,
             size=68, bold=True, line_spacing=1.0)
    add_text(slide, "live.", 1.1, 3.5, 11, 1.4,
             size=68, bold=True, color=TEAL, line_spacing=1.0)
    add_text(slide,
             "Exécution du pipeline · API en route · dashboard interactif.",
             1.1, 5.2, 11, 0.6, size=18, color=GREY)

    add_footer(slide, n, total)


def slide_team(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, NAVY)
    add_kicker(slide, "10 · Équipe & repo", GOLD)

    add_text(slide, "Merci.", 0.7, 1.2, 12, 1.6,
             size=88, bold=True, line_spacing=1.0)
    add_text(slide, "Questions, retours, démo approfondie.",
             0.7, 3.2, 12, 0.7, size=22, color=GREY)

    # team cards
    team = [
        ("Adam Beloucif",   "@Adam-Blf",   GOLD),
        ("Emilien Morice",  "@emilien754", TEAL),
    ]
    for i, (name, handle, color) in enumerate(team):
        x = 0.7 + i * 6.2
        add_round_rect(slide, x, 4.5, 5.9, 1.4, NAVY2)
        add_round_rect(slide, x, 4.5, 0.15, 1.4, color)
        add_text(slide, name,   x + 0.4, 4.65, 5.5, 0.5, size=22, bold=True, color=WHITE)
        add_text(slide, handle, x + 0.4, 5.15, 5.5, 0.4, size=14, color=color, font="Courier New")
        add_text(slide, "Data Engineer · M1 EFREI", x + 0.4, 5.45, 5.5, 0.4, size=11, color=GREY)

    add_text(slide,
             "github.com/Adam-Blf/urban-data-explorer",
             0.7, 6.2, 12, 0.5, size=14, color=TEAL, font="Courier New")

    add_footer(slide, n, total)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = Inches(WIDTH_IN)
    prs.slide_height = Inches(HEIGHT_IN)

    builders = [
        slide_cover,
        slide_problem,
        slide_answer,
        slide_architecture,
        slide_sources,
        slide_indicators,
        slide_pipeline,
        slide_api,
        slide_dashboard,
        slide_choices,
        slide_demo,
        slide_team,
    ]
    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        if i == 1:
            builder(prs)
        else:
            builder(prs, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"PPTX generated: {OUT} ({total} slides)")


if __name__ == "__main__":
    main()
